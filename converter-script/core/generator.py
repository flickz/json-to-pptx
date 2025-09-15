from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import List, Optional, Tuple
import os
import logging

from .data_parser import DataParser
from .models import Frame, TextWidget, ImageWidget
from .coordinate_converter import CoordinateConverter
from .text_extractor import ContentExtractor, ExtractedContent
from .image_handler import ImageHandler

logger = logging.getLogger(__name__)


class PowerPointGenerator:
    """Generates PowerPoint presentations from Miro data"""
    
    def __init__(self, slide_width_inches: float = 10.0, slide_height_inches: float = 5.625, 
                 image_cache_dir: str = ".image_cache", content_extractor: Optional[ContentExtractor] = None):
        """
        Initialize the PowerPoint generator
        
        Args:
            slide_width_inches: Width of the slide in inches
            slide_height_inches: Height of the slide in inches
            image_cache_dir: Directory for caching downloaded images
            content_extractor: Instance of ContentExtractor
        """
        self.slide_width_inches = slide_width_inches
        self.slide_height_inches = slide_height_inches
        self.presentation = None
        self.slide = None
        self.image_handler = ImageHandler(cache_dir=image_cache_dir)
        self.content_extractor = content_extractor or ContentExtractor()
        
        
    def create_presentation(self):
        """Create a new presentation with custom slide size"""
        self.presentation = Presentation()
        
        self.presentation.slide_width = Inches(self.slide_width_inches)
        self.presentation.slide_height = Inches(self.slide_height_inches)
        
        blank_slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        self.slide = self.presentation.slides.add_slide(blank_slide_layout)
        
    def set_slide_background(self, hex_color: str):
        """Set the slide background color"""
        # Convert hex to RGB
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        
        # Set background
        background = self.slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
        
    def hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def add_text_box(self, left_inches: float, top_inches: float, 
                     width_inches: float, height_inches: float,
                     content: ExtractedContent, font_size_pt: int):
        """
        Add a text box to the slide with content and styling
        
        Args:
            left_inches: Left position in inches
            top_inches: Top position in inches
            width_inches: Width in inches
            height_inches: Height in inches
            content: Extracted content with text and styling
            font_size_pt: Calculated font size in points
        """
        # Ensure minimum height for visibility
        MIN_HEIGHT = 0.3  # Minimum height in inches
        if height_inches < MIN_HEIGHT:
            # Adjust position to keep text centered
            top_inches = top_inches - (MIN_HEIGHT - height_inches) / 2
            height_inches = MIN_HEIGHT
        
        left = Inches(left_inches)
        top = Inches(top_inches)
        width = Inches(width_inches)
        height = Inches(height_inches)
        
        text_box = self.slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        
        if height_inches < 0.5:
            text_frame.margin_left = Inches(0.05)
            text_frame.margin_right = Inches(0.05)
            text_frame.margin_top = Inches(0.02)
            text_frame.margin_bottom = Inches(0.02)
        else:
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)
        
        text_frame.word_wrap = True
        
        if content.background_color:
            fill = text_box.fill
            fill.solid()
            r, g, b = self.hex_to_rgb(content.background_color)
            fill.fore_color.rgb = RGBColor(r, g, b)
        
        text_frame.clear()  # Clear any default text
        paragraph = text_frame.paragraphs[0]
        
        if content.text_align == 'center':
            paragraph.alignment = PP_ALIGN.CENTER
        elif content.text_align == 'right':
            paragraph.alignment = PP_ALIGN.RIGHT
        else:
            paragraph.alignment = PP_ALIGN.LEFT
        
        if content.text_align == 'center':
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        if content.processed_text.runs:
            for i, run_data in enumerate(content.processed_text.runs):
                if i == 0:
                    if paragraph.runs:
                        run = paragraph.runs[0]
                        run.text = run_data.text
                    else:
                        run = paragraph.add_run()
                        run.text = run_data.text
                else:
                    run = paragraph.add_run()
                    run.text = run_data.text
                
                # Set font
                font = run.font
                font.name = self.content_extractor.get_pptx_font_mapping(content.font_family)
                font.size = Pt(font_size_pt)
                
                # Set color
                if run_data.color:
                    r, g, b = self.hex_to_rgb(run_data.color)
                else:
                    r, g, b = self.hex_to_rgb(content.text_color)
                font.color.rgb = RGBColor(r, g, b)
                
                # Set formatting
                font.bold = run_data.bold
                font.italic = run_data.italic
                font.underline = run_data.underline
        else:
            if paragraph.runs:
                run = paragraph.runs[0]
                run.text = content.plain_text
            else:
                run = paragraph.add_run()
                run.text = content.plain_text
            
            font = run.font
            font.name = self.content_extractor.get_pptx_font_mapping(content.font_family)
            font.size = Pt(font_size_pt)
            
            r, g, b = self.hex_to_rgb(content.text_color)
            font.color.rgb = RGBColor(r, g, b)
            
            if content.base_style:
                font.bold = content.base_style.get('bold', False)
                font.italic = content.base_style.get('italic', False)
                font.underline = content.base_style.get('underline', False)
    
    def add_image(self, left_inches: float, top_inches: float,
                  width_inches: float, height_inches: float,
                  image_widget: ImageWidget):
        """
        Add an image to the slide
        
        Args:
            left_inches: Left position in inches
            top_inches: Top position in inches
            width_inches: Width in inches
            height_inches: Height in inches
            image_widget: ImageWidget with image data
        """
        if not image_widget.image_url:
            logger.warning(f"Warning: No image URL for widget {image_widget.id}")
            return
        
        image_path = self.image_handler.download_image(image_widget.image_url)
        if not image_path:
            logger.error(f"Error: Failed to download image from {image_widget.image_url}")
            return
        
        if (image_widget.crop and 
            image_widget.crop.width > 0 and 
            image_widget.crop.height > 0 and
            (image_widget.crop.x > 0 or 
             image_widget.crop.y > 0 or
             image_widget.crop.width != image_widget.resource.width or
             image_widget.crop.height != image_widget.resource.height)):
            
            cropped_path = self.image_handler.crop_image(
                image_path,
                image_widget.crop.x,
                image_widget.crop.y,
                image_widget.crop.width,
                image_widget.crop.height
            )
            if cropped_path:
                image_path = cropped_path
        
        try:
            left = Inches(left_inches)
            top = Inches(top_inches)
            width = Inches(width_inches)
            height = Inches(height_inches)
            
            pic = self.slide.shapes.add_picture(str(image_path), left, top, width, height)
            
            
            logger.info(f"  ✓ Added image: {image_widget.title or 'Untitled'}")
            logger.info(f"    Position: ({left_inches:.2f}\", {top_inches:.2f}\")")
            logger.info(f"    Size: {width_inches:.2f}\" x {height_inches:.2f}\"")
            
        except Exception as e:
            logger.error(f"Error adding image to slide: {str(e)}")
    
    def generate_from_json_data(self, data_path: str, output_path: str = "output.pptx"):
        """
        Generate PowerPoint from Miro data file
        
        Args:
            data_path: Path to the Miro data JSON file
            output_path: Path for the output PowerPoint file
        """
        logger.info(f"Generating PowerPoint from {data_path}")
        logger.info("=" * 60)
        
        parser = DataParser(data_path)
        widgets = parser.parse()
        
        frames = [w for w in widgets.values() if isinstance(w, Frame)]
        
        if not frames:
            raise ValueError("No frames found in the data")
        
        frames.sort(key=lambda f: f.presentation_order or "")
        
        logger.info(f"Found {len(frames)} frame(s) to process")
        
        self.presentation = Presentation()
        
        first_frame = frames[0]
        self.presentation.slide_width = Inches(self.slide_width_inches)
        self.presentation.slide_height = Inches(self.slide_height_inches)
        
        for frame_idx, frame in enumerate(frames):
            logger.info(f"\n{'='*60}")
            logger.info(f"Processing Frame {frame_idx + 1} of {len(frames)}")
            logger.info(f"Frame name: {frame.name or 'Unnamed'}")
            logger.info(f"Frame size: {frame.size.width:.2f} x {frame.size.height:.2f} px")
            logger.info(f"Presentation order: {frame.presentation_order or 'None'}")
            
            # Initialize coordinate converter for this frame
            coord_converter = CoordinateConverter(frame.size.width, frame.size.height,
                                                 self.slide_width_inches, self.slide_height_inches)
            
            blank_slide_layout = self.presentation.slide_layouts[6]  # Blank layout
            self.slide = self.presentation.slides.add_slide(blank_slide_layout)
            
            if frame.style and frame.style.background_color is not None:
                bg_color = self.content_extractor.int_to_hex_color(frame.style.background_color)
                if bg_color:
                    self.set_slide_background(bg_color)
                    logger.info(f"Slide background set to {bg_color}")
            
            frame_bounds = coord_converter.get_frame_bounds()
            
            logger.info("\nAdding elements to slide:")
            for i, child in enumerate(frame.children):
                # Each widget knows how to render itself
                child.render(self.slide, self, coord_converter, frame_bounds)
        
        self.presentation.save(output_path)

        logger.info(f"\n{'='*60}")
        logger.info(f"✅ Presentation saved to: {output_path}")
        logger.info(f"📊 Total slides created: {len(frames)}")
        
        return output_path
